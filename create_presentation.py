import pandas as pd
from pptx import Presentation
from pptx.util import Inches

def load_data(file_path):
    # Load the Excel file
    return pd.read_excel(file_path)

def create_presentation(data):
    # Initialize a presentation object
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the layout for a title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "2024 ABC Conference Agenda"
    subtitle.text = "Welcome to the Conference!"

    # Slide for each speaker
    for _, row in data.iterrows():
        slide_layout = prs.slide_layouts[1]  # 1 is the layout for a title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        # Setting the title and content for each speaker
        speaker_info = f"Speaker: {row['Speaker']}\nTitle & Affiliation: {row['Title & Affiliation']}"
        title.text = f"{row['Time']} - {row['Program/Section']}"
        content.text = speaker_info

        # Add agenda if available
        if pd.notna(row['Min']):
            content.text += f"\nDuration: {int(row['Min'])} minutes"

        # Adding LinkedIn link if available
        if pd.notna(row['LinkedIn']):
            content.text += f"\nLinkedIn: {row['LinkedIn']}"

    # Save the presentation
    prs.save('Conference_Agenda_Presentation.pptx')
    print("Presentation created successfully!")

def main():
    file_path = '2024 ABC Agenda.xlsx'
    data = load_data(file_path)
    create_presentation(data)

if __name__ == "__main__":
    main()
